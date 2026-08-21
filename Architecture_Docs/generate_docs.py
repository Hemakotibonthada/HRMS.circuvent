#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
generate_docs.py — Multi-format documentation builder for Career.circuvent.

Reads the numbered Markdown sources in this directory and produces:

    Architecture_Guide.md       aggregated single-page master reference
    Architecture_Guide.docx     styled Word document
    Architecture_Guide.pdf      cover page, table of contents, vector diagrams
    Architecture_Overview.pptx  14-slide presentation deck

Self-contained: no network access, no build step. Run it from anywhere.

    python generate_docs.py

Dependencies: python-docx, python-pptx, reportlab, markdown
"""

from __future__ import annotations

import os
import re
import sys
import datetime
from dataclasses import dataclass, field
from typing import Iterable

# ───────────────────────────────────────────────────────────── constants ──

HERE = os.path.dirname(os.path.abspath(__file__))

SOURCES = [
    ("01_SYSTEM_OVERVIEW.md", "System Overview"),
    ("02_DATABASE_AND_DATA_MODELS.md", "Database & Data Models"),
    ("03_INTEGRATIONS_AND_ECOSYSTEM.md", "Integrations & Ecosystem"),
    ("04_MAINTENANCE_AND_OPERATIONS.md", "Maintenance & Operations"),
    ("05_AREAS_OF_ENHANCEMENT.md", "Areas of Enhancement"),
    ("06_ARCHITECTURE_DIAGRAMS.md", "Architecture Diagram Atlas"),
]

TITLE = "hrms.circuvent.com"
SUBTITLE = "Architecture & Technical Audit"
ORG = "Circuvent Technologies"
BUILT = datetime.date.today().isoformat()

# Brand palette (hex, no leading #)
INK = "1F2933"
MUTED = "5A6875"
ACCENT = "1D4ED8"
ACCENT_DK = "1E3A8A"
GOOD = "15803D"
WARN = "B45309"
BAD = "B91C1C"
RULE = "D5DBE1"
CODE_BG = "F4F6F8"
BAND = "0F172A"

# Glyphs that no installed font renders, mapped per-context. Verified against
# the actual cmap tables of Consolas and Calibri with fontTools, not guessed.
#
#   mono  — used inside code blocks. MUST be a single character, because these
#           diagrams are column-aligned and a two-character substitution shears
#           every box below it.
#   prose — used in body text, tables, headings and slides, where a word is
#           clearer and width does not matter.
#
# Consolas has the full box-drawing and block-element range, so those survive
# in code. Calibri does not, so they are flattened in prose.
GLYPHS = {
    # (character): (mono replacement, prose replacement)
    "\u2705": ("Y", "YES"),          # ✅
    "\u274c": ("N", "NO"),           # ❌
    "\u26a0": ("!", "!"),            # ⚠
    "\u2713": ("+", "yes"),          # ✓
    "\u2717": ("x", "no"),           # ✗
    "\u2605": ("*", "*"),            # ★
    "\u2606": (".", "-"),            # ☆
    "\u25b6": (">", ">"),            # ▶
    "\u25c0": ("<", "<"),            # ◀
    "\u21d2": ("=>", "=>"),        # ⇒
    "\U0001f512": ("#", "[encrypted]"),  # 🔒
    "\U0001f534": ("!", "[CRITICAL]"),  # 🔴
    "\U0001f7e0": ("!", "[HIGH]"),   # 🟠
    "\U0001f7e1": ("~", "[MEDIUM]"),  # 🟡
    "\U0001f7e2": (".", "[LOW]"),    # 🟢
    "\u2699": ("*", ""),             # ⚙
    "\U0001f310": ("*", ""),         # 🌐
    "\U0001f4e6": ("*", "[pkg]"),    # 📦
    "\ufe0f": ("", ""),              # variation selector
    "\U0001f535": (".", "[INFO]"),  # 🔵
    "\u2b50": ("*", "*"),           # ⭐
    "\u2716": ("x", "x"),           # ✖
    "\u27e8": ("<", "<"),           # ⟨
    "\u27e9": (">", ">"),           # ⟩
    "\u2208": ("in", "in"),         # ∈
}

# Present in Consolas but absent from Calibri: flattened in prose only.
PROSE_ONLY = {
    "\u2550": "-", "\u2551": "|", "\u2588": "#", "\u2591": ":",
    "\u25bc": "v", "\u25b2": "^", "\u2570": "'", "\u256f": "'",
    "\u2554": "+", "\u2557": "+", "\u255a": "+", "\u255d": "+",
    "\u2560": "+", "\u2563": "+", "\u2564": "+",
}


def strip_glyphs(text: str, mono: bool = False) -> str:
    """Replace characters the target font cannot draw.

    `mono=True` keeps box-drawing intact (Consolas has it) and uses
    single-character substitutions so diagram columns stay aligned.
    """
    idx = 0 if mono else 1
    for ch, pair in GLYPHS.items():
        if ch in text:
            text = text.replace(ch, pair[idx])
    if not mono:
        for ch, rep in PROSE_ONLY.items():
            if ch in text:
                text = text.replace(ch, rep)
    return text


# ─────────────────────────────────────────────────────────── md parsing ──

@dataclass
class Block:
    kind: str                       # heading|para|code|table|list|quote|hr
    text: str = ""
    level: int = 0
    lang: str = ""
    lines: list = field(default_factory=list)
    rows: list = field(default_factory=list)
    ordered: bool = False


def _split_row(line: str) -> list:
    line = line.strip()
    if line.startswith("|"):
        line = line[1:]
    if line.endswith("|"):
        line = line[:-1]
    return [c.strip() for c in line.split("|")]


def parse_markdown(text: str) -> list:
    """Line-based Markdown parser covering the subset used by these documents."""
    out: list = []
    lines = text.replace("\r\n", "\n").split("\n")
    i = 0
    n = len(lines)

    while i < n:
        line = lines[i]
        stripped = line.strip()

        # fenced code
        if stripped.startswith("```"):
            lang = stripped[3:].strip()
            i += 1
            buf = []
            while i < n and not lines[i].strip().startswith("```"):
                buf.append(lines[i])
                i += 1
            i += 1
            while buf and not buf[-1].strip():
                buf.pop()
            out.append(Block("code", lang=lang, lines=buf))
            continue

        # blank
        if not stripped:
            i += 1
            continue

        # horizontal rule
        if re.fullmatch(r"(-{3,}|\*{3,}|_{3,})", stripped):
            out.append(Block("hr"))
            i += 1
            continue

        # heading
        m = re.match(r"^(#{1,6})\s+(.*)$", stripped)
        if m:
            out.append(Block("heading", text=m.group(2).strip(), level=len(m.group(1))))
            i += 1
            continue

        # table: a header row followed by a separator row
        if stripped.startswith("|") and i + 1 < n and re.match(
            r"^\s*\|?[\s:\-|]+\|[\s:\-|]*$", lines[i + 1]
        ) and "-" in lines[i + 1]:
            header = _split_row(lines[i])
            i += 2
            rows = []
            while i < n and lines[i].strip().startswith("|"):
                rows.append(_split_row(lines[i]))
                i += 1
            width = len(header)
            norm = []
            for r in rows:
                r = (r + [""] * width)[:width]
                norm.append(r)
            out.append(Block("table", rows=[header] + norm))
            continue

        # blockquote
        if stripped.startswith(">"):
            buf = []
            while i < n and lines[i].strip().startswith(">"):
                buf.append(re.sub(r"^\s*>\s?", "", lines[i]))
                i += 1
            merged = []
            cur = []
            for b in buf:
                if b.strip():
                    cur.append(b.strip())
                else:
                    if cur:
                        merged.append(" ".join(cur))
                        cur = []
            if cur:
                merged.append(" ".join(cur))
            out.append(Block("quote", lines=merged))
            continue

        # list
        if re.match(r"^\s*([-*+]|\d+\.)\s+", line):
            ordered = bool(re.match(r"^\s*\d+\.\s+", line))
            items = []
            while i < n and re.match(r"^\s*([-*+]|\d+\.)\s+", lines[i]):
                raw = lines[i]
                indent = len(raw) - len(raw.lstrip())
                body = re.sub(r"^\s*([-*+]|\d+\.)\s+", "", raw).strip()
                i += 1
                # continuation lines
                while (i < n and lines[i].strip()
                       and not re.match(r"^\s*([-*+]|\d+\.)\s+", lines[i])
                       and not lines[i].strip().startswith(("#", "|", "```", ">"))
                       and (len(lines[i]) - len(lines[i].lstrip())) > indent):
                    body += " " + lines[i].strip()
                    i += 1
                items.append((min(indent // 2, 2), body))
            out.append(Block("list", rows=items, ordered=ordered))
            continue

        # paragraph
        buf = []
        while (i < n and lines[i].strip()
               and not lines[i].strip().startswith(("#", "```", ">", "|"))
               and not re.match(r"^\s*([-*+]|\d+\.)\s+", lines[i])
               and not re.fullmatch(r"(-{3,}|\*{3,}|_{3,})", lines[i].strip())):
            buf.append(lines[i].strip())
            i += 1
        if buf:
            out.append(Block("para", text=" ".join(buf)))
        else:
            i += 1

    return out


# inline: -> list of (text, bold, italic, code, href)
_INLINE = re.compile(
    r"(?P<link>\[[^\]]+\]\([^)]+\))"
    r"|(?P<code>`[^`]+`)"
    r"|(?P<bold>\*\*[^*]+\*\*)"
    r"|(?P<italic>(?<!\*)\*[^*\n]+\*(?!\*))"
)


def inline_tokens(text: str):
    pos = 0
    for m in _INLINE.finditer(text):
        if m.start() > pos:
            yield (text[pos:m.start()], False, False, False, None)
        raw = m.group(0)
        if m.group("code"):
            yield (raw[1:-1], False, False, True, None)
        elif m.group("link"):
            lm = re.match(r"\[([^\]]+)\]\(([^)]+)\)", raw)
            label = lm.group(1).strip()
            if label.startswith("`") and label.endswith("`"):
                label = label[1:-1]
            yield (label, False, False, False, lm.group(2))
        elif m.group("bold"):
            # recurse: bold spans may wrap code, links or italics
            for t, _b, i, c, h in inline_tokens(raw[2:-2]):
                yield (t, True, i, c, h)
        else:
            for t, b, _i, c, h in inline_tokens(raw[1:-1]):
                yield (t, b, True, c, h)
        pos = m.end()
    if pos < len(text):
        yield (text[pos:], False, False, False, None)


# ────────────────────────────────────────────── 1 · master markdown ──

def build_master_md() -> str:
    parts = [
        f"# {TITLE} — {SUBTITLE}\n",
        f"> **Organisation:** {ORG}  ",
        f"> **Generated:** {BUILT}  ",
        "> **Scope:** full technical audit and architecture reverse-engineering.\n",
        "\nThis is the aggregated master reference. The same content is maintained "
        "as five focused documents in this directory; edit those, then re-run "
        "`generate_docs.py` to rebuild this file and the Word, PDF and PowerPoint "
        "deliverables.\n",
        "\n---\n",
        "\n## Contents\n",
    ]
    for idx, (_, label) in enumerate(SOURCES, start=1):
        anchor = "part-%d-%s" % (idx, re.sub(r"[^a-z0-9]+", "-", label.lower()).strip("-"))
        parts.append(f"{idx}. [Part {idx} · {label}](#{anchor})")
    parts.append("\n---\n")

    for idx, (fname, label) in enumerate(SOURCES, start=1):
        with open(os.path.join(HERE, fname), "r", encoding="utf-8") as fh:
            body = fh.read()
        # cross-document links become plain text inside a single file
        body = re.sub(r"\[`?([^\]`]+)`?\]\(\./\d\d_[A-Z_]+\.md\)", r"**\1**", body)
        body = re.sub(r"^#\s+.*$", "", body, count=1, flags=re.M).lstrip("\n")
        parts.append(f"\n<a id=\"part-{idx}-"
                     f"{re.sub(r'[^a-z0-9]+', '-', label.lower()).strip('-')}\"></a>\n")
        parts.append(f"\n# Part {idx} · {label}\n")
        parts.append(body.rstrip())
        parts.append("\n\n---\n")

    parts.append(f"\n*Generated by `generate_docs.py` on {BUILT}.*\n")
    return "\n".join(parts)


# ────────────────────────────────────────────────────── 2 · docx ──

def build_docx(path: str) -> None:
    from docx import Document
    from docx.shared import Pt, Inches, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
    from docx.enum.table import WD_TABLE_ALIGNMENT
    from docx.enum.section import WD_SECTION
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    sec = doc.sections[0]
    sec.left_margin = Inches(0.85)
    sec.right_margin = Inches(0.85)
    sec.top_margin = Inches(0.85)
    sec.bottom_margin = Inches(0.8)

    normal = doc.styles["Normal"]
    normal.font.name = "Calibri"
    normal.font.size = Pt(10.5)
    normal.font.color.rgb = RGBColor.from_string(INK)
    normal.paragraph_format.space_after = Pt(6)
    normal.paragraph_format.line_spacing = 1.15
    rpr = normal.element.get_or_add_rPr()
    rf = rpr.find(qn("w:rFonts"))
    if rf is None:
        rf = OxmlElement("w:rFonts")
        rpr.append(rf)
    rf.set(qn("w:eastAsia"), "Calibri")

    for name, size, color, bold in (
        ("Heading 1", 20, ACCENT_DK, True),
        ("Heading 2", 15, ACCENT_DK, True),
        ("Heading 3", 12.5, INK, True),
        ("Heading 4", 11, MUTED, True),
        ("Heading 5", 10.5, MUTED, True),
        ("Heading 6", 10, MUTED, True),
    ):
        st = doc.styles[name]
        st.font.name = "Calibri"
        st.font.size = Pt(size)
        st.font.bold = bold
        st.font.color.rgb = RGBColor.from_string(color)
        st.paragraph_format.space_before = Pt(14 if name == "Heading 1" else 10)
        st.paragraph_format.space_after = Pt(6)
        st.paragraph_format.keep_with_next = True

    def shade(paragraph, hexcolor):
        pr = paragraph._p.get_or_add_pPr()
        sh = OxmlElement("w:shd")
        sh.set(qn("w:val"), "clear")
        sh.set(qn("w:fill"), hexcolor)
        pr.append(sh)

    def left_bar(paragraph, hexcolor):
        pr = paragraph._p.get_or_add_pPr()
        bd = OxmlElement("w:pBdr")
        lf = OxmlElement("w:left")
        lf.set(qn("w:val"), "single")
        lf.set(qn("w:sz"), "18")
        lf.set(qn("w:space"), "8")
        lf.set(qn("w:color"), hexcolor)
        bd.append(lf)
        pr.append(bd)

    def add_runs(paragraph, text, base_size=10.5, mono=False):
        text = strip_glyphs(text)
        for chunk, bold, italic, code, href in inline_tokens(text):
            if not chunk:
                continue
            run = paragraph.add_run(chunk)
            run.bold = bold
            run.italic = italic
            if code or mono:
                run.font.name = "Consolas"
                run.font.size = Pt(base_size - 1)
                run.font.color.rgb = RGBColor.from_string(BAD if code else INK)
                r = run._element.get_or_add_rPr().get_or_add_rFonts()
                r.set(qn("w:ascii"), "Consolas")
                r.set(qn("w:hAnsi"), "Consolas")
            else:
                run.font.size = Pt(base_size)
            if href:
                run.font.color.rgb = RGBColor.from_string(ACCENT)
                run.underline = True

    # ── cover ──
    for _ in range(5):
        doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(TITLE)
    r.font.size = Pt(34)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string(ACCENT_DK)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run(SUBTITLE)
    r.font.size = Pt(17)
    r.font.color.rgb = RGBColor.from_string(MUTED)

    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = p.add_run("─" * 30)
    r.font.color.rgb = RGBColor.from_string(RULE)

    for line, size, color in (
        (ORG, 12, INK),
        ("Multi-tenant HRMS · Next.js 16 · Neon Postgres · Drizzle · Vercel", 11, MUTED),
        (f"Generated {BUILT}", 10, MUTED),
    ):
        p = doc.add_paragraph()
        p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(line)
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor.from_string(color)

    doc.add_paragraph()
    facts = doc.add_table(rows=0, cols=2)
    facts.style = "Light List Accent 1"
    facts.alignment = WD_TABLE_ALIGNMENT.CENTER
    for k, v in (
        ("Architecture", "Multi-tenant HRMS — the system of record for who works here"),
        ("Scale", "1,454 files · 574 TypeScript · 144,146 lines · 150 API routes · 102 pages"),
        ("Database", "Neon Postgres · 117 tables · 44 enums · 39 migrations · FORCE RLS"),
        ("Identity", "Four sign-in paths converging on one HS256 suite session"),
        ("Tests", "92 files · 2,664 passing · plus nine custom verification scripts"),
        ("CI/CD", "The only real CI in the suite — running 9 of its own 12 checks"),
        ("Mobile", "Kotlin Multiplatform, shipped at v1.8.0 / versionCode 10"),
        ("Primary risk", "CI has no live-database check — the blind spot that already failed once"),
    ):
        row = facts.add_row().cells
        run = row[0].paragraphs[0].add_run(k)
        run.bold = True
        run.font.size = Pt(9.5)
        run2 = row[1].paragraphs[0].add_run(v)
        run2.font.size = Pt(9.5)

    doc.add_page_break()

    # ── table of contents field ──
    h = doc.add_paragraph()
    r = h.add_run("Contents")
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = RGBColor.from_string(ACCENT_DK)

    p = doc.add_paragraph()
    fld = OxmlElement("w:fldSimple")
    fld.set(qn("w:instr"), r'TOC \o "1-3" \h \z \u')
    inner = OxmlElement("w:r")
    t = OxmlElement("w:t")
    t.text = "Right-click and choose “Update Field” to build the table of contents."
    inner.append(t)
    fld.append(inner)
    p._p.append(fld)

    doc.add_page_break()

    # ── body ──
    for idx, (fname, label) in enumerate(SOURCES, start=1):
        with open(os.path.join(HERE, fname), "r", encoding="utf-8") as fh:
            blocks = parse_markdown(fh.read())

        if idx > 1:
            doc.add_page_break()

        band = doc.add_paragraph()
        shade(band, BAND)
        r = band.add_run(f"  PART {idx}")
        r.font.size = Pt(9)
        r.font.bold = True
        r.font.color.rgb = RGBColor.from_string("FFFFFF")

        first_heading_used = False

        for b in blocks:
            if b.kind == "heading":
                if not first_heading_used and b.level == 1:
                    first_heading_used = True
                    doc.add_heading(strip_glyphs(b.text), level=1)
                    continue
                doc.add_heading(strip_glyphs(b.text), level=min(b.level, 6))

            elif b.kind == "para":
                p = doc.add_paragraph()
                add_runs(p, b.text)

            elif b.kind == "quote":
                for ln in b.lines:
                    p = doc.add_paragraph()
                    p.paragraph_format.left_indent = Inches(0.22)
                    p.paragraph_format.space_before = Pt(4)
                    left_bar(p, ACCENT)
                    shade(p, "EEF3FB")
                    add_runs(p, ln, base_size=10)
                    for run in p.runs:
                        run.italic = True

            elif b.kind == "code":
                is_mermaid = b.lang.lower() == "mermaid"
                if is_mermaid:
                    cap = doc.add_paragraph()
                    cr = cap.add_run("Mermaid diagram — renders in the Markdown edition")
                    cr.font.size = Pt(8)
                    cr.bold = True
                    cr.font.color.rgb = RGBColor.from_string(ACCENT)
                    cap.paragraph_format.space_after = Pt(0)
                widest = max((len(x) for x in b.lines), default=0)
                size = 8 if widest <= 92 else (7 if widest <= 108 else 6)
                for ln in b.lines:
                    p = doc.add_paragraph()
                    p.paragraph_format.space_after = Pt(0)
                    p.paragraph_format.space_before = Pt(0)
                    p.paragraph_format.line_spacing = 1.0
                    shade(p, "EDF2F7" if is_mermaid else CODE_BG)
                    run = p.add_run(strip_glyphs(ln, mono=True) or " ")
                    run.font.name = "Consolas"
                    run.font.size = Pt(size)
                    run.font.color.rgb = RGBColor.from_string(INK)
                    rf2 = run._element.get_or_add_rPr().get_or_add_rFonts()
                    rf2.set(qn("w:ascii"), "Consolas")
                    rf2.set(qn("w:hAnsi"), "Consolas")
                doc.add_paragraph().paragraph_format.space_after = Pt(2)

            elif b.kind == "table":
                header, *rows = b.rows
                has_header = any(c.strip() for c in header)
                if not has_header:
                    rows = [header] + rows
                ncols = len(b.rows[0])
                t = doc.add_table(rows=0, cols=ncols)
                t.style = "Light Grid Accent 1" if has_header else "Light List Accent 1"
                t.autofit = True
                if has_header:
                    cells = t.add_row().cells
                    for c, txt in zip(cells, header):
                        para = c.paragraphs[0]
                        add_runs(para, txt, base_size=9)
                        for run in para.runs:
                            run.bold = True
                for row in rows:
                    cells = t.add_row().cells
                    for c, txt in zip(cells, row):
                        para = c.paragraphs[0]
                        para.paragraph_format.space_after = Pt(2)
                        add_runs(para, txt, base_size=9)
                doc.add_paragraph().paragraph_format.space_after = Pt(2)

            elif b.kind == "list":
                for depth, item in b.rows:
                    style = "List Number" if b.ordered else "List Bullet"
                    try:
                        p = doc.add_paragraph(style=style)
                    except KeyError:
                        p = doc.add_paragraph()
                    p.paragraph_format.left_indent = Inches(0.25 + 0.22 * depth)
                    p.paragraph_format.space_after = Pt(2)
                    add_runs(p, item, base_size=10)

            elif b.kind == "hr":
                p = doc.add_paragraph()
                r = p.add_run("─" * 58)
                r.font.color.rgb = RGBColor.from_string(RULE)
                r.font.size = Pt(8)

    doc.save(path)


# ─────────────────────────────────────────────────────── 3 · pdf ──

def build_pdf(path: str) -> None:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm, mm
    from reportlab.lib import colors
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.enums import TA_CENTER, TA_LEFT
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont
    from reportlab.platypus import (
        BaseDocTemplate, PageTemplate, Frame, Paragraph, Spacer, Table,
        TableStyle, PageBreak, KeepTogether, Flowable,
    )
    from reportlab.platypus.tableofcontents import TableOfContents

    fonts = r"C:\Windows\Fonts"

    def reg(name, filename, fallback):
        try:
            pdfmetrics.registerFont(TTFont(name, os.path.join(fonts, filename)))
            return name
        except Exception:
            return fallback

    BODY = reg("Body", "calibri.ttf", "Helvetica")
    BODY_B = reg("BodyB", "calibrib.ttf", "Helvetica-Bold")
    BODY_I = reg("BodyI", "calibrii.ttf", "Helvetica-Oblique")
    MONO = reg("Mono", "consola.ttf", "Courier")
    MONO_B = reg("MonoB", "consolab.ttf", "Courier-Bold")

    try:
        from reportlab.pdfbase.pdfmetrics import registerFontFamily
        registerFontFamily(BODY, normal=BODY, bold=BODY_B, italic=BODY_I, boldItalic=BODY_B)
        registerFontFamily(MONO, normal=MONO, bold=MONO_B, italic=MONO, boldItalic=MONO_B)
    except Exception:
        pass

    C_INK = colors.HexColor("#" + INK)
    C_MUTED = colors.HexColor("#" + MUTED)
    C_ACCENT = colors.HexColor("#" + ACCENT)
    C_ACCENT_DK = colors.HexColor("#" + ACCENT_DK)
    C_RULE = colors.HexColor("#" + RULE)
    C_CODE = colors.HexColor("#" + CODE_BG)
    C_BAND = colors.HexColor("#" + BAND)

    PW, PH = A4
    LM = RM = 1.5 * cm
    TM = 1.7 * cm
    BM = 1.6 * cm
    AVAIL = PW - LM - RM

    ss = getSampleStyleSheet()

    def st(name, **kw):
        base = dict(fontName=BODY, fontSize=9.5, leading=13.2, textColor=C_INK,
                    spaceBefore=0, spaceAfter=5, alignment=TA_LEFT)
        base.update(kw)
        return ParagraphStyle(name, **base)

    S_H1 = st("H1", fontName=BODY_B, fontSize=19, leading=23, textColor=C_ACCENT_DK,
              spaceBefore=6, spaceAfter=10)
    S_H2 = st("H2", fontName=BODY_B, fontSize=14, leading=18, textColor=C_ACCENT_DK,
              spaceBefore=13, spaceAfter=6)
    S_H3 = st("H3", fontName=BODY_B, fontSize=11.5, leading=15, textColor=C_INK,
              spaceBefore=10, spaceAfter=4)
    S_H4 = st("H4", fontName=BODY_B, fontSize=10, leading=13.5, textColor=C_MUTED,
              spaceBefore=8, spaceAfter=3)
    S_BODY = st("Body")
    S_QUOTE = st("Quote", fontName=BODY_I, fontSize=9.2, leading=13,
                 textColor=C_ACCENT_DK, leftIndent=9, spaceBefore=3, spaceAfter=6)
    S_LIST = st("List", leftIndent=12, spaceAfter=2.5)
    S_CELL = st("Cell", fontSize=8.2, leading=10.6, spaceAfter=0)
    S_CELLH = st("CellH", fontName=BODY_B, fontSize=8.2, leading=10.6,
                 spaceAfter=0, textColor=colors.white)
    S_CAP = st("Cap", fontName=BODY_B, fontSize=7.6, leading=9.6,
               textColor=C_ACCENT, spaceAfter=1.5)
    S_COVER_T = st("CoverT", fontName=BODY_B, fontSize=33, leading=38,
                   textColor=C_ACCENT_DK, alignment=TA_CENTER, spaceAfter=6)
    S_COVER_S = st("CoverS", fontSize=16, leading=21, textColor=C_MUTED,
                   alignment=TA_CENTER, spaceAfter=16)
    S_COVER_M = st("CoverM", fontSize=10.5, leading=15, textColor=C_MUTED,
                   alignment=TA_CENTER, spaceAfter=3)
    S_TOC1 = st("TOC1", fontName=BODY_B, fontSize=11, leading=18, spaceBefore=7)
    S_TOC2 = st("TOC2", fontSize=9.4, leading=14.5, leftIndent=16)
    S_TOC3 = st("TOC3", fontSize=8.6, leading=13, leftIndent=32, textColor=C_MUTED)

    def esc(text: str) -> str:
        return (text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;"))

    def rl_inline(text: str) -> str:
        text = strip_glyphs(text)
        out = []
        for chunk, bold, italic, code, href in inline_tokens(text):
            if not chunk:
                continue
            piece = esc(chunk)
            if code:
                piece = f'<font face="{MONO}" size="8.4" color="#{BAD}">{piece}</font>'
            if bold:
                piece = f"<b>{piece}</b>"
            if italic:
                piece = f"<i>{piece}</i>"
            if href:
                safe = href.replace("&", "&amp;").replace('"', "%22")
                piece = f'<link href="{safe}" color="#{ACCENT}">{piece}</link>'
            out.append(piece)
        return "".join(out) or "&nbsp;"

    # ── vector diagrams drawn natively ──
    class Diagram(Flowable):
        """Native vector drawing — sharp at any zoom, no raster assets."""

        def __init__(self, kind, width, height):
            super().__init__()
            self.kind = kind
            self.width = width
            self.height = height

        def wrap(self, aw, ah):
            return (self.width, self.height)

        def _box(self, x, y, w, h, label, sub=None, fill="#FFFFFF",
                 stroke="#1D4ED8", tcol="#1F2933", radius=5, bold=True):
            c = self.canv
            c.saveState()
            c.setFillColor(colors.HexColor(fill))
            c.setStrokeColor(colors.HexColor(stroke))
            c.setLineWidth(1.1)
            c.roundRect(x, y, w, h, radius, stroke=1, fill=1)
            c.setFillColor(colors.HexColor(tcol))
            c.setFont(BODY_B if bold else BODY, 8.4)
            if sub:
                c.drawCentredString(x + w / 2, y + h / 2 + 3.4, label)
                c.setFont(BODY, 7)
                c.setFillColor(colors.HexColor("#" + MUTED))
                c.drawCentredString(x + w / 2, y + h / 2 - 6.4, sub)
            else:
                c.drawCentredString(x + w / 2, y + h / 2 - 3, label)
            c.restoreState()

        def _arrow(self, x1, y1, x2, y2, label=None, col="#5A6875", dashed=False,
                   head=True):
            import math
            c = self.canv
            c.saveState()
            c.setStrokeColor(colors.HexColor(col))
            c.setLineWidth(1.0)
            if dashed:
                c.setDash(2.5, 2.5)
            c.line(x1, y1, x2, y2)
            c.setDash()
            if head:
                ang = math.atan2(y2 - y1, x2 - x1)
                spread = 0.40           # ~23 degrees each side
                size = 5.2
                c.setFillColor(colors.HexColor(col))
                p = c.beginPath()
                p.moveTo(x2, y2)
                p.lineTo(x2 - size * math.cos(ang - spread),
                         y2 - size * math.sin(ang - spread))
                p.lineTo(x2 - size * math.cos(ang + spread),
                         y2 - size * math.sin(ang + spread))
                p.close()
                c.drawPath(p, stroke=0, fill=1)
            if label:
                c.setFillColor(colors.HexColor(col))
                c.setFont(BODY, 6.6)
                c.drawCentredString((x1 + x2) / 2, (y1 + y2) / 2 + 3.6, label)
            c.restoreState()

        def draw(self):
            getattr(self, "_draw_" + self.kind)()

        def _draw_context(self):
            W, H = self.width, self.height
            cy = H - 112
            bw, bh = W * 0.56, 34
            x_app = (W - bw) / 2

            up = [("auth.circuvent.com", "suite OIDC provider \u00b7 RS256 / JWKS",
                   "#F5F3FF", "#6D28D9"),
                  ("Customer IdPs \u00b7 SCIM 2.0", "per-tenant Okta \u00b7 Entra \u00b7 Google",
                   "#F8FAFC", "#64748B")]
            uw = (W - 90) / 2
            up_y = H - 34
            u_centres = []
            for n, (lab, sub, fill, stroke) in enumerate(up):
                x = n * (uw + 90)
                self._box(x, up_y, uw, 28, lab, sub, fill=fill, stroke=stroke)
                u_centres.append(x + uw / 2)

            bus_y = up_y - 14
            c = self.canv
            c.saveState()
            c.setStrokeColor(colors.HexColor("#64748B"))
            c.setLineWidth(0.9)
            c.line(u_centres[0], bus_y, u_centres[-1], bus_y)
            for cx in u_centres:
                c.line(cx, up_y, cx, bus_y)
            c.restoreState()
            self._arrow(W / 2, bus_y, W / 2, cy + bh + 4,
                        "relying party  \u00b7  inbound provisioning", col="#1D4ED8")

            self._box(x_app, cy, bw, bh, "hrms.circuvent.com",
                      "150 API routes \u00b7 102 pages \u00b7 117 tables",
                      fill="#DBEAFE", stroke="#1D4ED8")

            row_y = cy - 84
            deps = [("Neon Postgres", "identity + hrms \u00b7 FORCE RLS \u00b7 hrms_app",
                     "#EFF6FF", "#1D4ED8"),
                    ("Cloudflare R2", "signed PDFs \u00b7 fails HARD",
                     "#ECFDF5", "#15803D"),
                    ("paystub.circuvent", "employee master \u00b7 outbox",
                     "#FEF3C7", "#B45309")]
            gap = 8.0
            dw = (W - gap * (len(deps) - 1)) / len(deps)
            d_centres = []
            for n, (lab, sub, fill, stroke) in enumerate(deps):
                x = n * (dw + gap)
                self._box(x, row_y, dw, 34, lab, sub, fill=fill, stroke=stroke)
                d_centres.append((x + dw / 2, stroke))

            dbus = cy - 22
            c = self.canv
            c.saveState()
            c.setStrokeColor(colors.HexColor("#1D4ED8"))
            c.setLineWidth(1.0)
            c.line(W / 2, cy - 3, W / 2, dbus)
            c.line(d_centres[0][0], dbus, d_centres[-1][0], dbus)
            c.restoreState()
            for cx, col in d_centres:
                self._arrow(cx, dbus, cx, row_y + 34 + 4, col=col)

            c = self.canv
            c.saveState()
            c.setFont(BODY, 6.6)
            c.setFillColor(colors.HexColor("#" + MUTED))
            c.drawCentredString(W / 2, row_y - 13,
                                "The suite shares ONE Postgres project and reads the same "
                                "identity schema directly. That, not any API, is the real "
                                "backbone of single sign-on.")
            c.restoreState()

        def _draw_layers(self):
            W, H = self.width, self.height
            rows = [
                ("THE GATE", "middleware.ts - verifies, then OVERWRITES x-user-id / x-org-id / x-user-role",
                 "#F5F3FF", "#6D28D9"),
                ("ROUTES", "150 handlers - 128 session, 3 API key, 2 SCIM, 17 public or bespoke",
                 "#EFF6FF", "#1D4ED8"),
                ("AUTHORIZATION", "rbac.ts - ~90 permissions, 4 roles + owner - fails closed",
                 "#ECFDF5", "#15803D"),
                ("PURE RULE MODULES", "~40 modules - statutory-india, rostering, workflow, settlement",
                 "#FFF7ED", "#B45309"),
                ("REPOSITORIES", "*.neon.ts - the only code that queries - always via withTenant()",
                 "#F0FDF4", "#15803D"),
                ("POSTGRES", "117 tables - FORCE ROW LEVEL SECURITY - the role must NOT bypass it",
                 "#F1F5F9", "#5A6875"),
            ]
            h = 26
            gap = 6
            y = H - h
            for lab, sub, fill, stroke in rows:
                self._box(0, y, W, h, "", None, fill=fill, stroke=stroke)
                c = self.canv
                c.saveState()
                c.setFillColor(colors.HexColor(stroke))
                c.setFont(BODY_B, 8.0)
                c.drawString(9, y + h / 2 + 2.4, lab)
                c.setFillColor(colors.HexColor("#" + MUTED))
                c.setFont(BODY, 6.9)
                c.drawString(9, y + h / 2 - 7.4, sub)
                c.restoreState()
                y -= (h + gap)

            c = self.canv
            c.saveState()
            c.setFont(BODY, 6.6)
            c.setFillColor(colors.HexColor("#" + MUTED))
            c.drawCentredString(W / 2, y + h - 4,
                                "Nine verification scripts guard this stack. The two that "
                                "check the live deployment are not in CI.")
            c.restoreState()

        def _draw_apply(self):
            W, H = self.width, self.height
            lanes = ["Browser", "middleware", "Route", "Postgres", "Effects"]
            n = len(lanes)
            lw = W / n
            c = self.canv
            top = H - 16
            for i, lab in enumerate(lanes):
                x = i * lw + lw / 2
                c.saveState()
                c.setFillColor(colors.HexColor("#" + ACCENT_DK))
                c.setFont(BODY_B, 7.4)
                c.drawCentredString(x, top, lab)
                c.setStrokeColor(colors.HexColor("#" + RULE))
                c.setLineWidth(0.8)
                c.setDash(2, 3)
                c.line(x, top - 8, x, 10)
                c.setDash()
                c.restoreState()
            steps = [
                (0, 1, "GET /api/employees  -  cv_access cookie"),
                (1, 1, "public prefix? no. jwtVerify - NO database call"),
                (1, 0, "expired + refresh cookie -> 401 + x-session-refresh: 1"),
                (1, 1, "OVERWRITE x-user-id / x-org-id / x-user-role"),
                (1, 2, "forward"),
                (2, 2, "requireApiContext - orgId from the VERIFIED TOKEN"),
                (2, 2, "rbac - roleHasPermission / canAccessModule"),
                (2, 3, "withTenant(ctx) - set_config('app.org_id', ..., true)"),
                (3, 3, "RLS: org_id = app_current_org() - and the role must not bypass"),
                (3, 2, "rows for THIS tenant, or zero. Never another's."),
                (2, 4, "outbox row commits WITH the change - drained by /api/cron"),
                (2, 0, "the resource, or a typed error"),
            ]
            y = top - 22
            for a, b, lab in steps:
                xa = a * lw + lw / 2
                xb = b * lw + lw / 2
                if a == b:
                    c.saveState()
                    c.setStrokeColor(colors.HexColor("#" + WARN))
                    c.setLineWidth(1.0)
                    c.rect(xa, y - 4, 30, 10, stroke=1, fill=0)
                    c.setFillColor(colors.HexColor("#" + WARN))
                    c.setFont(BODY, 6.3)
                    c.drawString(xa + 34, y - 1, lab)
                    c.restoreState()
                else:
                    self._arrow(xa, y, xb, y, lab,
                                col="#1D4ED8" if b > a else "#15803D")
                y -= 15.6

    # ── flowable builders ──
    def code_flowable(block):
        # Measure the real rendered width instead of guessing from the
        # character count: box-drawing glyphs get substituted before
        # layout, and a hand-tuned per-character estimate under-measures
        # and silently clipped wide ASCII diagrams at the right margin.
        from reportlab.pdfbase.pdfmetrics import stringWidth
        measured = [strip_glyphs(x, mono=True) for x in block.lines] or [""]
        size = 7.6
        while size > 4.2:
            widest = max((stringWidth(x, MONO, size) for x in measured),
                         default=0.0)
            if widest <= AVAIL - 14:
                break
            size -= 0.2
        leading = size * 1.32
        is_mermaid = block.lang.lower() == "mermaid"

        frame_h = PH - TM - BM
        # Leave room for the caption, padding and a little slack, then work out
        # how many lines of this size fit on one page. A single un-splittable
        # table taller than the frame is a hard LayoutError, so long blocks —
        # the big Mermaid class diagram in particular — are chunked instead.
        per_page = max(8, int((frame_h - 60) / leading))

        chunks = [block.lines[i:i + per_page]
                  for i in range(0, len(block.lines), per_page)] or [[""]]

        parts = []
        for n, chunk in enumerate(chunks):
            body = "\n".join(strip_glyphs(x, mono=True) or " " for x in chunk)
            para = Paragraph(
                esc(body).replace("\n", "<br/>").replace(" ", "&nbsp;"),
                ParagraphStyle("Code", fontName=MONO, fontSize=size,
                               leading=leading, textColor=C_INK,
                               spaceBefore=0, spaceAfter=0),
            )
            t = Table([[para]], colWidths=[AVAIL])
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, -1),
                 colors.HexColor("#EDF2F7") if is_mermaid else C_CODE),
                ("BOX", (0, 0), (-1, -1), 0.6,
                 C_ACCENT if is_mermaid else C_RULE),
                ("LEFTPADDING", (0, 0), (-1, -1), 6),
                ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                ("TOPPADDING", (0, 0), (-1, -1), 5),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ]))
            if n == 0:
                if is_mermaid:
                    parts.append(Paragraph(
                        "Mermaid diagram &mdash; renders interactively in the "
                        "Markdown edition", S_CAP))
                parts.append(t)
            else:
                parts.append(PageBreak())
                parts.append(Paragraph(
                    f"(continued &mdash; part {n + 1} of {len(chunks)})", S_CAP))
                parts.append(t)
        parts.append(Spacer(1, 6))
        return parts

    def table_flowable(block):
        header, *rows = block.rows
        ncols = len(header)
        # A Markdown table may legitimately have an empty header row — several
        # in these documents are two-column fact lists. Rendering the dark
        # header band for one produces a blank blue bar, so drop it instead.
        has_header = any(c.strip() for c in header)
        if not has_header:
            header = []
        weights = []
        for i in range(ncols):
            col = ([header[i]] if has_header else []) + [r[i] for r in rows]
            weights.append(max(6, min(60, max((len(str(x)) for x in col),
                                              default=6))))
        total = float(sum(weights)) or 1.0
        widths = [max(38.0, AVAIL * w / total) for w in weights]
        scale = AVAIL / sum(widths)
        widths = [w * scale for w in widths]

        data = []
        if has_header:
            data.append([Paragraph(rl_inline(c), S_CELLH) for c in header])
        for r in rows:
            data.append([Paragraph(rl_inline(c), S_CELL) for c in r])
        if not data:
            return []

        t = Table(data, colWidths=widths,
                  repeatRows=1 if has_header else 0, hAlign="LEFT")
        style = [
            ("GRID", (0, 0), (-1, -1), 0.4, C_RULE),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ("LEFTPADDING", (0, 0), (-1, -1), 4.5),
            ("RIGHTPADDING", (0, 0), (-1, -1), 4.5),
            ("TOPPADDING", (0, 0), (-1, -1), 3.5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3.5),
        ]
        if has_header:
            style += [("BACKGROUND", (0, 0), (-1, 0), C_ACCENT_DK),
                      ("TEXTCOLOR", (0, 0), (-1, 0), colors.white)]
        else:
            style.append(("BACKGROUND", (0, 0), (0, -1),
                          colors.HexColor("#F1F5F9")))
        start = 1 if has_header else 0
        for i in range(start, len(data)):
            if (i - start) % 2 == 1:
                style.append(("BACKGROUND", (0, i), (-1, i),
                              colors.HexColor("#F8FAFC")))
        t.setStyle(TableStyle(style))
        return [t, Spacer(1, 7)]

    def quote_flowable(lines):
        inner = [Paragraph(rl_inline(x), S_QUOTE) for x in lines]
        t = Table([[inner]], colWidths=[AVAIL])
        t.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#EEF3FB")),
            ("LINEBEFORE", (0, 0), (0, -1), 2.4, C_ACCENT),
            ("LEFTPADDING", (0, 0), (-1, -1), 9),
            ("RIGHTPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
        ]))
        return [t, Spacer(1, 6)]

    # ── document scaffolding ──
    class Doc(BaseDocTemplate):
        def afterFlowable(self, flowable):
            if not isinstance(flowable, Paragraph):
                return
            style = flowable.style.name
            level = {"H1": 0, "H2": 1, "H3": 2}.get(style)
            if level is None:
                return
            text = re.sub(r"<[^>]+>", "", flowable.getPlainText())
            self.notify("TOCEntry", (level, text, self.page))

    def decorate(canv, doc_):
        canv.saveState()
        if doc_.page > 1:
            canv.setFillColor(C_BAND)
            canv.rect(0, PH - 13 * mm, PW, 13 * mm, stroke=0, fill=1)
            canv.setFillColor(colors.white)
            canv.setFont(BODY_B, 7.6)
            canv.drawString(LM, PH - 8.7 * mm, TITLE)
            canv.setFont(BODY, 7.6)
            canv.drawRightString(PW - RM, PH - 8.7 * mm, SUBTITLE)

            canv.setStrokeColor(C_RULE)
            canv.setLineWidth(0.6)
            canv.line(LM, BM - 4, PW - RM, BM - 4)
            canv.setFillColor(C_MUTED)
            canv.setFont(BODY, 7.2)
            canv.drawString(LM, BM - 13, f"{ORG} · generated {BUILT}")
            canv.drawRightString(PW - RM, BM - 13, f"Page {doc_.page}")
        canv.restoreState()

    doc = Doc(path, pagesize=A4, leftMargin=LM, rightMargin=RM,
              topMargin=TM, bottomMargin=BM,
              title=f"{TITLE} — {SUBTITLE}", author=ORG,
              subject="Architecture and technical audit")
    frame = Frame(LM, BM, AVAIL, PH - TM - BM, id="body",
                  leftPadding=0, rightPadding=0, topPadding=0, bottomPadding=0)
    doc.addPageTemplates([PageTemplate(id="main", frames=[frame], onPage=decorate)])

    story = []

    # cover
    story.append(Spacer(1, 3.6 * cm))
    story.append(Paragraph(TITLE, S_COVER_T))
    story.append(Paragraph(SUBTITLE, S_COVER_S))
    rule = Table([[""]], colWidths=[6 * cm], rowHeights=[2.2])
    rule.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), C_ACCENT)]))
    rule.hAlign = "CENTER"
    story.append(rule)
    story.append(Spacer(1, 0.8 * cm))
    story.append(Paragraph(ORG, S_COVER_M))
    story.append(Paragraph("Next.js 16 · Neon Postgres · Drizzle · Argon2id · SCIM 2.0 · Kotlin Multiplatform",
                           S_COVER_M))
    story.append(Paragraph(f"Generated {BUILT}", S_COVER_M))
    story.append(Spacer(1, 1.1 * cm))
    story.append(Diagram("context", AVAIL, 232))
    story.append(Spacer(1, 0.7 * cm))

    facts = [
        ["Architecture", "Multi-tenant HRMS — the system of record for who works here"],
        ["Scale", "1,454 files · 574 TypeScript · 144,146 lines · 150 API routes · 102 pages"],
        ["Database", "Neon Postgres · 117 tables · 44 enums · 39 migrations · FORCE RLS"],
        ["Identity", "Four sign-in paths converging on one HS256 suite session"],
        ["Tests", "92 files · 2,664 passing · plus nine custom verification scripts"],
        ["CI/CD", "The only real CI in the suite — running 9 of its own 12 checks"],
        ["Mobile", "Kotlin Multiplatform, shipped at v1.8.0 / versionCode 10"],
        ["Largest risk", "CI has no live-database check — the blind spot that already failed once"],
    ]
    ft = Table([[Paragraph(f"<b>{a}</b>", S_CELL), Paragraph(b, S_CELL)]
                for a, b in facts], colWidths=[AVAIL * 0.32, AVAIL * 0.68])
    ft.setStyle(TableStyle([
        ("GRID", (0, 0), (-1, -1), 0.4, C_RULE),
        ("BACKGROUND", (0, 0), (0, -1), colors.HexColor("#F1F5F9")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ("TOPPADDING", (0, 0), (-1, -1), 4),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
    ]))
    story.append(ft)
    story.append(PageBreak())

    # contents
    story.append(Paragraph("Contents", S_H1))
    toc = TableOfContents()
    toc.levelStyles = [S_TOC1, S_TOC2, S_TOC3]
    story.append(toc)
    story.append(PageBreak())

    # body
    for idx, (fname, label) in enumerate(SOURCES, start=1):
        with open(os.path.join(HERE, fname), "r", encoding="utf-8") as fh:
            blocks = parse_markdown(fh.read())

        if idx > 1:
            story.append(PageBreak())

        band = Table([[Paragraph(
            f'<font color="white" size="8"><b>PART {idx}</b></font>', S_CELL)]],
            colWidths=[AVAIL])
        band.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), C_BAND),
            ("LEFTPADDING", (0, 0), (-1, -1), 8),
            ("TOPPADDING", (0, 0), (-1, -1), 3),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
        ]))
        story.append(band)
        story.append(Spacer(1, 8))

        if idx == 1:
            pass

        for b in blocks:
            if b.kind == "heading":
                style = {1: S_H1, 2: S_H2, 3: S_H3}.get(b.level, S_H4)
                story.append(Paragraph(rl_inline(b.text), style))
                # insert the native diagrams alongside the sections they explain
                if idx == 1 and b.level == 2 and b.text.startswith("5."):
                    story.append(Spacer(1, 3))
                    story.append(Paragraph("Figure — layered architecture", S_CAP))
                    story.append(Diagram("layers", AVAIL, 200))
                    story.append(Spacer(1, 8))
                if idx == 1 and b.level == 3 and b.text.startswith("6.1"):
                    story.append(Spacer(1, 3))
                    story.append(Paragraph("Figure — one authenticated request, "
                                           "end to end", S_CAP))
                    story.append(Diagram("apply", AVAIL, 214))
                    story.append(Spacer(1, 8))

            elif b.kind == "para":
                story.append(Paragraph(rl_inline(b.text), S_BODY))

            elif b.kind == "quote":
                story.extend(quote_flowable(b.lines))

            elif b.kind == "code":
                story.extend(code_flowable(b))

            elif b.kind == "table":
                story.extend(table_flowable(b))

            elif b.kind == "list":
                counter = 0
                for depth, item in b.rows:
                    if b.ordered:
                        counter += 1
                        marker = f"<b>{counter}.</b>"
                    else:
                        marker = "&bull;"
                    story.append(Paragraph(
                        f"{marker}&nbsp;&nbsp;{rl_inline(item)}",
                        ParagraphStyle(f"L{depth}", parent=S_LIST,
                                       leftIndent=12 + 14 * depth)))
                story.append(Spacer(1, 4))

            elif b.kind == "hr":
                story.append(Spacer(1, 3))
                hr = Table([[""]], colWidths=[AVAIL], rowHeights=[0.7])
                hr.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), C_RULE)]))
                story.append(hr)
                story.append(Spacer(1, 6))

    doc.multiBuild(story)


# ────────────────────────────────────────────────────── 4 · pptx ──

SLIDES = [
    {
        "kind": "bullets",
        "title": "What this system is",
        "lead": "The system of record for who works here — and the largest "
                "application in the Circuvent suite.",
        "bullets": [
            ("A full multi-tenant HRMS: employees, leave, attendance, "
             "rostering, payroll, recruitment, performance, learning, "
             "documents and e-signature, helpdesk, assets, benefits, "
             "compensation, expenses and privacy governance.", 0),
            ("144,146 lines across 574 TypeScript source files, serving 150 "
             "API routes and 102 pages over 117 database tables.", 0),
            ("Plus a shipped native Android app — Kotlin Multiplatform, "
             "version 1.8.0, versionCode 10, signed and released.", 0),
            ("It owns its own payroll module. It does NOT delegate payroll "
             "to paystub.circuvent — the only link is a one-way push of "
             "employee master data.", 0),
            ("It is a relying party to two OIDC systems and a receiver of "
             "SCIM provisioning. It is never itself an identity provider.", 0),
        ],
    },
    {
        "kind": "facts",
        "title": "At a glance",
        "facts": [
            ("Stack", "Next.js 16.1 App Router (Turbopack) · React 19.2 · TS 5 strict"),
            ("Database", "Neon Postgres · Drizzle 0.45 · 117 tables · 44 enums"),
            ("Tenancy", "FORCE ROW LEVEL SECURITY · 117 tenant_isolation policies"),
            ("Passwords", "Argon2id — 19 MiB, t=2, p=1, PHC, auto-rehash"),
            ("Sessions", "HS256 · cv_access 15 min · cv_refresh 30 d, rotated"),
            ("Tests", "92 files · 2,664 passing · 39 migrations"),
            ("Guards", "Nine custom verification scripts + gitleaks"),
            ("CI/CD", "The only real CI in the suite — runs 9 of its own 12 checks"),
        ],
    },
    {
        "kind": "topology",
        'top': [('auth.circuvent.com', 'suite OIDC provider', 'violet'), ('hrms.circuvent.com', '150 routes · 102 pages', 'primary'), ('Customer IdPs', 'per-tenant OIDC + SCIM', 'neutral')],
        'links': ['OIDC · RS256', 'SCIM 2.0'],
        'hub': 1,
        'children': [('Neon Postgres', '117 tables · RLS · hrms_app'), ('Cloudflare R2', 'signed PDFs · fails HARD'), ('SMTP relay', 'nodemailer · fails SOFT'), ('paystub.circuvent', 'employee master · outbox')],
        'childnote': 'The suite shares ONE Postgres project. That, not any API, is the real backbone of single sign-on.',
        "title": "Architectural topology",
        "note": "ATS, Mail and DevOps appear only as nav links. No direct REST calls to them exist.",
    },
    {
        "kind": "bullets",
        "title": "The finding that defines this codebase",
        "lead": "From the header of scripts/smoke-live.ts.",
        "bullets": [
            ("\u201cninety-one correct policies and seventy-five passing "
             "isolation tests, while DATABASE_URL pointed at a role with "
             "BYPASSRLS and every query returned every tenant's rows. "
             "Nothing that ran in CI could have noticed.\u201d", 0),
            ("The policies were right. The tests were right. The tests "
             "passed. And every tenant could read every other tenant's data.", 0),
            ("The role hrms_app existed with rolbypassrls = false — but was "
             "never granted LOGIN. So DATABASE_URL fell back to "
             "neondb_owner, the table owner, which bypasses RLS regardless.", 0),
            ("Fixed by migration 0028. Guarded at runtime by "
             "assertConnectionIsolatesTenants(), which refuses to start "
             "against a bypassing role.", 0),
            ("The lesson is not \u201cadd a test\u201d. It is that some "
             "properties are only true of a DEPLOYMENT, and a test suite "
             "cannot see them.", 0),
        ],
    },
    {
        "kind": "layers",
        "title": "Layered design",
        "note": "One gate, ~40 pure rule modules, and isolation delegated entirely to Postgres.",
    },
    {
        "kind": "flow",
        "title": "Data flow - one authenticated request",
        "steps": [
            ("1", "Gate", "middleware.ts - public prefix? no. jwtVerify, no DB call."),
            ("2", "Refresh", "expired + refresh cookie -> 401 + x-session-refresh: 1"),
            ("3", "Headers", "OVERWRITE x-user-id / x-org-id / x-user-role"),
            ("4", "Context", "requireApiContext - orgId from the VERIFIED TOKEN"),
            ("5", "Permit", "rbac.ts - ~90 permissions - canAccessModule fails closed"),
            ("6", "Tenant", "withTenant(ctx) - set_config('app.org_id', ..., true)"),
            ("7", "Postgres", "RLS: org_id = app_current_org(), role must NOT bypass"),
            ("8", "Effects", "outbox row commits WITH the change, drained by /api/cron"),
        ],
    },
    {
        "kind": "bullets",
        "title": "Four sign-in paths, one session",
        "lead": "They converge, and none of them can skip MFA.",
        "bullets": [
            ("Local password — Argon2id, 19 MiB / t=2 / p=1, PHC-encoded, "
             "auto-rehashed when the parameters drift.", 0),
            ("Suite OIDC — relying party to auth.circuvent.com, with PKCE, "
             "state and nonce, verified against published JWKS.", 0),
            ("Per-tenant customer OIDC — each tenant's own Okta, Entra or "
             "Google, selected by email domain. Explicitly NOT SAML, with "
             "the XML signature-wrapping attack class named as the reason.", 0),
            ("WebAuthn passkeys — credential_id unique globally.", 0),
            ("All four mint the same HS256 token: sub, orgId, role, email. "
             "cv_access 15 minutes; cv_refresh 30 days, SHA-256 hashed, "
             "single-use, rotated — and REUSE REVOKES THE WHOLE FAMILY.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "MFA, done carefully",
        "lead": "A three-state machine, and every state earns its existence.",
        "bullets": [
            ("off -> pending -> active. The secret is stored ENCRYPTED "
             "(AES-256-GCM) from the moment it is minted.", 0),
            ("Backup codes are issued ONLY after a live code proves the "
             "secret works. Codes are never handed out for an unproven "
             "secret.", 0),
            ("`pending` is deliberately not enforced at sign-in — anti-"
             "lockout — and it grants no elevated trust either.", 0),
            ("Disabling requires the current password AND a live TOTP code. "
             "A stolen session cookie alone cannot turn the control off.", 0),
            ("VERIFIED: both signIn() and signInWithSso() independently call "
             "mfaRequiredAtSignIn(). There is no bypass.", 0),
            ("The trade is fail-closed: an MFA-enabled user currently cannot "
             "complete SSO or passkey login at all. An availability gap, "
             "not a security gap.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "Tenancy: no WHERE clause anywhere",
        "lead": "Isolation is delegated to Postgres, so application code cannot forget it.",
        "bullets": [
            ("116 of 117 tables carry org_id. The exception is "
             "`organizations` — it IS the tenant.", 0),
            ("One reusable function, apply_tenant_rls(), sweeps every table "
             "with an org_id column and is called by 17 later migrations "
             "rather than being hand-repeated 117 times.", 0),
            ("FORCE ROW LEVEL SECURITY, not merely ENABLE — so the policy "
             "binds the table OWNER too: \u201ca mistake in a migration "
             "script must not be able to read across tenants either.\u201d", 0),
            ("withTenant(ctx, fn) is the only sanctioned entry point. It "
             "sets app.org_id with SET LOCAL semantics, so the value reverts "
             "before the connection returns to the pool.", 0),
            ("Forgetting withTenant() is FAIL-CLOSED: app_current_org() "
             "returns NULL, the predicate is false, zero rows come back — a "
             "visible bug, not a silent leak.", 0),
        ],
    },
    {
        "kind": "facts",
        "title": "The schema, in one table",
        "facts": [
            ("Total", "117 physical tables — 116 in Drizzle, 1 raw SQL only"),
            ("Schemas", "identity (20, cross-app) · hrms (97, HR domain)"),
            ("Enums", "44 types, from asset_state to signature_status"),
            ("Money", "*_minor bigint. Never numeric. Never float."),
            ("Deletes", "CASCADE or SET NULL, with exactly three RESTRICT"),
            ("Audit", "Hash-chained AND append-only AND grant-revoked"),
            ("doc_store", "Schemaless jsonb catch-all for ~20 legacy collections"),
            ("Caching", "None. No Redis, no materialised views. All reads hit Postgres."),
        ],
    },
    {
        "kind": "bullets",
        "title": "Nine guard scripts that tests cannot replace",
        "lead": "Each one checks a property a unit test structurally cannot see.",
        "bullets": [
            ("audit-data-paths — reconciles 91 pages against their data "
             "sources. A silent 404 rendering as \u201cno items yet\u201d is "
             "indistinguishable from genuinely having no data.", 0),
            ("audit-fabricated-data — hardcoded values that look live. An "
             "earlier version scanned 2 of 9 directories and reported clean.", 0),
            ("verify-modules — 65 checks, because \u201ctwo routes were once "
             "fakes that returned 201 and wrote nothing.\u201d", 0),
            ("verify-query-plans — seeds 4,000 rows, EXPLAINs, then DROPS "
             "THE INDEX to prove the index is what removed the sort.", 0),
            ("verify-live-isolation — connects as whoever we really connect "
             "as, plants a row in one tenant and reads as another.", 0),
            ("verify-credential-reach — proves each app's credential cannot "
             "open a sibling app's database.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "The blind spot that remains",
        "lead": "CI is real. It is also structurally unable to catch what already happened.",
        "bullets": [
            (".github/workflows/verify.yml runs typecheck, lint:strict, four "
             "database verifies, two audits, 2,664 tests, a build, and a "
             "gitleaks history scan. That is more than any sibling app has.", 0),
            ("But every db:verify* step runs against IN-MEMORY PGlite, which "
             "connects as its own superuser. RLS policy correctness passes "
             "regardless of which role production actually uses.", 0),
            ("The two scripts that WOULD have caught the incident — "
             "db:verify:live and db:verify:reach — need real credentials, "
             "and are therefore the two not in CI.", 0),
            ("CI also omits three of the twelve checks in `npm run verify`: "
             "typecheck:mobile, lint:a11y and audit:unwired. A developer "
             "running verify locally is STRICTER than the merge gate.", 0),
            ("And verify-live-isolation's core cross-tenant test has never "
             "actually executed — only one organisation exists, so it is "
             "skipped, not passed.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "Money: mostly right, with one honest seam",
        "lead": "The rule is stated in the codebase's own words.",
        "bullets": [
            ("money/minor.ts: \u201cA whole number of paise, as a decimal "
             "string. Never a float.\u201d Held throughout "
             "statutory-india.ts, compensation.ts, settlement.ts, assets.ts "
             "and expense-rules.ts.", 0),
            ("But payroll-engine.ts is 580 lines of plain `number` — a "
             "legacy remnant whose five other major functions have zero "
             "callers.", 0),
            ("Two of them are still wired into the live pipeline: "
             "BigInt(Math.round(calculateProfessionalTax(minorToMajor(gross)) "
             "* 100)) — bigint, to float, back to bigint.", 0),
            ("minorToMajor() is the codebase's own DISPLAY-ONLY helper, "
             "whose doc comment says the result \u201cmust never be summed "
             "or compared for equality.\u201d", 0),
            ("And gratuity exists three times — one correct bigint "
             "implementation, two naive float duplicates with no part-year "
             "rounding and no death waiver. Both dead. Both still exported.", 0),
        ],
    },
    {
        "kind": "bullets",
        "title": "The same rules, in three codebases",
        "lead": "A statutory rate change must be made three times, in two languages.",
        "bullets": [
            ("src/lib/statutory-india.ts — web TypeScript. Ships.", 0),
            ("android/shared/ — Kotlin Multiplatform. Ships, at v1.8.0.", 0),
            ("mobile/src/lib/ — Expo TypeScript. Abandoned; its own "
             "documentation says it \u201chas never been run on a "
             "device\u201d.", 0),
            ("All three declare the same application id: "
             "com.circuvent.hrms. Nothing checks that they agree.", 0),
            ("A naming trap follows from it: leave-rules.ts, shift-rules.ts "
             "and attendance-rules.ts do NOT exist in src/lib. Searching "
             "those names finds only the abandoned copy.", 0),
        ],
    },
    {
        "kind": "scorecard",
        "title": "Health assessment",
        "rows": [
            ("Verification discipline", 5, "Nine guard scripts; real CI"),
            ("Test coverage", 4, "2,664 tests across 92 files"),
            ("Code hygiene", 4, "0 TODO, 0 ts-ignore, 0 console.log in src/"),
            ("Domain rigour", 4, "Pure modules; dated statutory config"),
            ("Auth & session design", 4, "Argon2id; refresh family revocation"),
            ("Tenancy enforcement", 4, "Token-derived orgId in 149 of 150 routes"),
            ("Money consistency", 3, "One float seam; three gratuities"),
            ("CI completeness", 3, "9 of 12 checks; no live-database check"),
            ("Documentation accuracy", 2, "Two docs describe a dead system"),
            ("Encryption at rest", 2, "bank_details is plaintext jsonb"),
            ("Rule-source singularity", 2, "Business rules live in three codebases"),
            ("Response headers", 1, "next.config.ts has none at all"),
            ("Observability", 1, "No APM, no error tracking, no alerting"),
        ],
    },
    {
        "kind": "bullets",
        "title": "Top five gaps",
        "lead": "Four are defects. The first is the capacity to notice the next one.",
        "bullets": [
            ("CI has no live-database check. Every DB verify runs against "
             "PGlite as its own superuser — the exact blind spot that let "
             "the BYPASSRLS incident through.  [1 week]", 0),
            ("/api/documents/reminders reads orgId from the QUERY STRING "
             "after a static shared token, with no rate limit. The one "
             "tenancy exception in 150 routes.  [half a day]", 0),
            ("next.config.ts has no security headers at all — no CSP, HSTS, "
             "X-Frame-Options or Referrer-Policy — on an app holding salary "
             "and bank details.  [1 hour]", 0),
            ("bank_details is unencrypted jsonb. Full account number and "
             "IFSC in plaintext; masked only on read.  [3 weeks]", 0),
            ("Two migrations are missing from the journal and would never "
             "run in a fresh environment. This is currently the one failing "
             "check in verify-migrations.ts.  [10 minutes]", 0),
        ],
    },
    {
        "kind": "roadmap",
        "title": "Phased roadmap",
        "phases": [
            ("Phase 1", "about 1 week", "Close what is already open",
             "Add the two missing journal entries · security headers in "
             "next.config.ts · derive orgId from the token on "
             "documents/reminders · delete mobile/ · rewrite the two "
             "obsolete architecture documents"),
            ("Phase 2", "2-3 weeks", "Make CI tell the truth",
             "A CI job with a real two-organisation Neon branch running "
             "db:verify:live and db:verify:reach · add the three omitted "
             "verify checks · seed a second org so the cross-tenant test "
             "finally executes"),
            ("Phase 3", "about 1 month", "Finish the money migration",
             "Port professional tax and new-regime income tax to bigint · "
             "delete payroll-engine.ts entirely · delete the two duplicate "
             "gratuity implementations · assert there is only one"),
            ("Phase 4-6", "1-2 quarters", "Data at rest, domain gaps, operations",
             "Encrypt bank_details via a column type migration · monthly "
             "leave accrual and balance validation · delete 1,583 lines of "
             "dead modules · error tracking, structured logs and outbox "
             "alerting · durable rate limiting · extend lint:strict past "
             "its allow-list"),
        ],
    },
    {
        "kind": "bullets",
        "title": "What must NOT change",
        "lead": "Decisions that look like overhead and are not.",
        "bullets": [
            ("The pure-core / impure-shell split. It is why 2,664 tests run "
             "in under a minute with no database.", 0),
            ("assertConnectionIsolatesTenants() — the single best line of "
             "defence in the codebase.", 0),
            ("FORCE ROW LEVEL SECURITY, and the one function that sweeps it "
             "across every table rather than 117 hand-written policies.", 0),
            ("The transactional outbox, all four times. The queue row and "
             "the business change commit together, or neither does.", 0),
            ("MFA that will not issue backup codes for an unproven secret, "
             "and interview_scorecards.submitted_at gating visibility — "
             "anchoring bias prevented in the data model, not a UI rule.", 0),
            ("The comment convention where a module opens by naming the "
             "specific bug it exists to prevent. It made this audit "
             "materially easier, and will make every future one easier too.", 0),
        ],
    },
    {
        "kind": "closing",
        "title": "In one sentence",
        "lead": "A system that learned the hard lesson, and has not quite finished acting on it.",
        "body": "One hundred and forty-four thousand lines with no TODO "
                "markers, no @ts-ignore, no suppressed lint rules and no "
                "stray console.log. Nine verification scripts that check "
                "things a test suite structurally cannot. Argon2id with "
                "auto-rehash, refresh-token family revocation, a hash-chained "
                "append-only audit log, and orgId derived from a verified "
                "token in 149 of 150 routes.\n\n"
                "All of that exists because of one incident: ninety-one "
                "correct policies, seventy-five passing tests, and a "
                "connection string that made every one of them inert. The "
                "codebase responded by building the only real CI in the "
                "suite \u2014 and the two checks that would have caught it "
                "are still the two that are not in it.\n\n"
                "Close that gap, finish the bigint migration, encrypt "
                "bank_details, add security headers, and delete the two "
                "architecture documents that describe a system which no "
                "longer exists. Then this is the strongest application in "
                "the suite by a clear margin.",
    },
]


def build_pptx(path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SW = prs.slide_width
    SH = prs.slide_height
    blank = prs.slide_layouts[6]

    def rgb(h):
        return RGBColor.from_string(h)

    def textbox(slide, x, y, w, h, text, size=16, bold=False, color=INK,
                align=PP_ALIGN.LEFT, font="Calibri", anchor=MSO_ANCHOR.TOP,
                spacing=1.0):
        tb = slide.shapes.add_textbox(x, y, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = anchor
        lines = text.split("\n")
        for i, ln in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = spacing
            r = p.add_run()
            r.text = strip_glyphs(ln)
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = font
            r.font.color.rgb = rgb(color)
        return tb

    def rect(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
        sh = slide.shapes.add_shape(shape, x, y, w, h)
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill)
        if line:
            sh.line.color.rgb = rgb(line)
            sh.line.width = Pt(1.2)
        else:
            sh.line.fill.background()
        sh.shadow.inherit = False
        sh.text_frame.word_wrap = True
        return sh

    def label(shape, text, size=12, bold=True, color="FFFFFF",
              align=PP_ALIGN.CENTER):
        tf = shape.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for i, ln in enumerate(text.split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = strip_glyphs(ln)
            r.font.size = Pt(size)
            r.font.bold = bold and i == 0
            r.font.name = "Calibri"
            r.font.color.rgb = rgb(color)

    def chrome(slide, title, index):
        rect(slide, 0, 0, SW, Inches(0.92), BAND, shape=MSO_SHAPE.RECTANGLE)
        textbox(slide, Inches(0.55), Inches(0.16), SW - Inches(2.4), Inches(0.6),
                title, size=25, bold=True, color="FFFFFF")
        rect(slide, 0, Inches(0.92), SW, Emu(38100), ACCENT,
             shape=MSO_SHAPE.RECTANGLE)
        textbox(slide, SW - Inches(1.5), Inches(0.26), Inches(1.0), Inches(0.4),
                f"{index:02d}", size=15, bold=True, color="64748B",
                align=PP_ALIGN.RIGHT)
        textbox(slide, Inches(0.55), SH - Inches(0.46),
                Inches(8.0), Inches(0.3),
                f"{TITLE} · {SUBTITLE} · {BUILT}", size=9, color=MUTED)

    for i, spec in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank)
        kind = spec["kind"]

        # ── title slide ──
        if kind == "title":
            rect(slide, 0, 0, SW, SH, BAND, shape=MSO_SHAPE.RECTANGLE)
            rect(slide, 0, SH - Inches(0.13), SW, Inches(0.13), ACCENT,
                 shape=MSO_SHAPE.RECTANGLE)
            textbox(slide, Inches(1.1), Inches(2.35), SW - Inches(2.2), Inches(1.3),
                    spec["title"], size=54, bold=True, color="FFFFFF",
                    align=PP_ALIGN.CENTER)
            rl = rect(slide, (SW - Inches(3.0)) // 2, Inches(3.72), Inches(3.0),
                      Emu(38100), ACCENT, shape=MSO_SHAPE.RECTANGLE)
            textbox(slide, Inches(1.1), Inches(4.05), SW - Inches(2.2), Inches(1.2),
                    spec["subtitle"], size=19, color="C7D2DE",
                    align=PP_ALIGN.CENTER, spacing=1.35)
            textbox(slide, Inches(1.1), Inches(5.72), SW - Inches(2.2), Inches(0.9),
                    "151,500 lines  ·  112 API routes  ·  "
                    "one 956 MB mail server", size=13, color="7F8EA3",
                    align=PP_ALIGN.CENTER)
            continue

        chrome(slide, spec["title"], i)
        top = Inches(1.28)

        # ── bulleted slide ──
        if kind == "bullets":
            y = top
            if spec.get("lead"):
                bar = rect(slide, Inches(0.55), y, SW - Inches(1.1), Inches(0.62),
                           "EEF3FB", ACCENT)
                label(bar, spec["lead"], size=16, color=ACCENT_DK,
                      align=PP_ALIGN.LEFT)
                bar.text_frame.margin_left = Inches(0.22)
                y += Inches(0.88)
            for text, depth in spec["bullets"]:
                h = Inches(0.52) if len(text) < 96 else Inches(0.78)
                x = Inches(0.75) + Inches(0.5) * depth
                w = SW - Inches(1.5) - Inches(0.5) * depth
                dot = rect(slide, x, y + Inches(0.13), Inches(0.12), Inches(0.12),
                           ACCENT if depth == 0 else "94A3B8", shape=MSO_SHAPE.OVAL)
                textbox(slide, x + Inches(0.3), y, w - Inches(0.3), h, text,
                        size=15 if depth == 0 else 13.5,
                        color=INK if depth == 0 else MUTED, spacing=1.12)
                y += h + Inches(0.06)

        # ── fact table ──
        elif kind == "facts":
            y = top + Inches(0.08)
            rowh = Inches(0.62)
            for n, (k, v) in enumerate(spec["facts"]):
                bg = "F1F5F9" if n % 2 == 0 else "FFFFFF"
                rect(slide, Inches(0.55), y, SW - Inches(1.1), rowh, bg,
                     shape=MSO_SHAPE.RECTANGLE)
                rect(slide, Inches(0.55), y, Emu(38100), rowh, ACCENT,
                     shape=MSO_SHAPE.RECTANGLE)
                textbox(slide, Inches(0.78), y + Inches(0.13), Inches(3.5),
                        Inches(0.4), k, size=14, bold=True, color=ACCENT_DK)
                textbox(slide, Inches(4.45), y + Inches(0.13),
                        SW - Inches(5.1), Inches(0.4), v, size=14, color=INK)
                y += rowh + Inches(0.07)

        # ── topology diagram ──
        elif kind == "topology":
            # Data-driven. spec keys:
            #   top       [(label, sub, tone)]  tone: neutral | primary | warn
            #   links     [str] connectors between adjacent top boxes
            #   hub       index of the box the children hang beneath
            #   children  [(label, sub)] up to 5
            #   childnote caption under the children row
            TONES = {
                "neutral": ("F8FAFC", "64748B", INK),
                "primary": ("DBEAFE", ACCENT, ACCENT_DK),
                "warn": ("FEF3C7", WARN, "7C2D12"),
                "good": ("ECFDF5", GOOD, "14532D"),
                "violet": ("F5F3FF", "6D28D9", "4C1D95"),
            }
            top = spec.get("top") or [("Clients", "browser", "neutral"),
                                      (TITLE, "this application", "primary")]
            links = spec.get("links") or []
            hub = spec.get("hub", len(top) - 1)
            children = spec.get("children") or []

            n_top = len(top)
            gap = Inches(1.15) if n_top > 2 else Inches(1.6)
            bw = min(Inches(2.9), (SW - Inches(1.1) - gap * (n_top - 1)) // n_top)
            bh = Inches(1.15)
            span = bw * n_top + gap * (n_top - 1)
            x0 = (SW - span) // 2
            cy = Inches(1.7)

            xs = [x0 + i * (bw + gap) for i in range(n_top)]
            for i, (lab, sub, tone) in enumerate(top):
                fill, stroke, ink = TONES.get(tone, TONES["neutral"])
                bx = rect(slide, xs[i], cy, bw, bh, fill, stroke)
                label(bx, f"{lab}\n{sub}" if sub else lab, size=14, color=ink)

            for i, txt in enumerate(links[: n_top - 1]):
                xa = xs[i] + bw
                rect(slide, xa + Inches(0.1), cy + bh // 2 - Emu(19050),
                     gap - Inches(0.2), Emu(38100), "94A3B8",
                     shape=MSO_SHAPE.RECTANGLE)
                textbox(slide, xa, cy + bh // 2 - Inches(0.44), gap, Inches(0.32),
                        txt, size=10, color=MUTED, align=PP_ALIGN.CENTER)

            if children:
                row_y = cy + bh + Inches(1.05)
                row_h = Inches(0.95)
                bus_y = cy + bh + Inches(0.52)
                sgap = Inches(0.22)
                margin = Inches(0.55)
                sw = (SW - margin * 2 - sgap * (len(children) - 1)) // len(children)
                centres = []
                for n, (lab, sub) in enumerate(children):
                    x = margin + n * (sw + sgap)
                    bx = rect(slide, x, row_y, sw, row_h, "ECFDF5", GOOD)
                    label(bx, f"{lab}\n{sub}" if sub else lab, size=12,
                          color="14532D")
                    centres.append(x + sw // 2)

                hub_cx = xs[max(0, min(hub, n_top - 1))] + bw // 2
                rect(slide, hub_cx - Emu(19050), cy + bh, Emu(38100),
                     bus_y - (cy + bh), GOOD, shape=MSO_SHAPE.RECTANGLE)
                rect(slide, centres[0], bus_y, centres[-1] - centres[0],
                     Emu(38100), GOOD, shape=MSO_SHAPE.RECTANGLE)
                for cx in centres:
                    rect(slide, cx - Emu(19050), bus_y, Emu(38100),
                         row_y - bus_y, GOOD, shape=MSO_SHAPE.RECTANGLE)

                cnote = spec.get("childnote")
                if cnote:
                    textbox(slide, margin, row_y + row_h + Inches(0.16),
                            SW - margin * 2, Inches(0.34), cnote, size=11,
                            color="14532D", align=PP_ALIGN.CENTER)

            if spec.get("note"):
                textbox(slide, Inches(0.6), SH - Inches(1.5),
                        SW - Inches(1.2), Inches(0.5), spec["note"],
                        size=13, bold=True, color=ACCENT_DK,
                        align=PP_ALIGN.CENTER)

        # ── layer stack ──
        elif kind == "layers":
            rows = [
                ("PRESENTATION", "React 19 server + client components · Tailwind v4",
                 "EFF6FF", ACCENT),
                ("APPLICATION", "App Router pages · server actions · route handlers",
                 "F5F3FF", "6D28D9"),
                ("DOMAIN — pure, no I/O",
                 "filter · format · resume · registration-validate · portal-groups",
                 "ECFDF5", GOOD),
                ("INTEGRATION", "ats.ts · candidate.ts · insights-reporter.ts",
                 "FFF7ED", WARN),
                ("PLATFORM", "Node.js on Vercel · CDN · 60-second ISR cache",
                 "F1F5F9", "64748B"),
            ]
            y = top + Inches(0.1)
            h = Inches(0.86)
            for lab, sub, fill, stroke in rows:
                bx = rect(slide, Inches(0.75), y, SW - Inches(1.5), h, fill, stroke)
                tf = bx.text_frame
                tf.margin_left = Inches(0.28)
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                r = p.add_run()
                r.text = lab
                r.font.size = Pt(15)
                r.font.bold = True
                r.font.name = "Calibri"
                r.font.color.rgb = rgb(stroke)
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.LEFT
                r2 = p2.add_run()
                r2.text = strip_glyphs(sub)
                r2.font.size = Pt(11.5)
                r2.font.name = "Calibri"
                r2.font.color.rgb = rgb(MUTED)
                y += h + Inches(0.14)
            textbox(slide, Inches(0.75), SH - Inches(0.95), SW - Inches(1.5),
                    Inches(0.4), spec["note"], size=12, color=ACCENT_DK,
                    align=PP_ALIGN.CENTER)

        # ── numbered flow ──
        elif kind == "flow":
            steps = spec["steps"]
            cols = 4
            rows_n = 2
            cw = (SW - Inches(1.5) - Inches(0.28) * (cols - 1)) / cols
            ch = Inches(1.95)
            for n, (num, head, body) in enumerate(steps):
                r_i, c_i = divmod(n, cols)
                x = Inches(0.75) + c_i * (cw + Inches(0.28))
                y = top + Inches(0.25) + r_i * (ch + Inches(0.42))
                card = rect(slide, x, y, cw, ch, "FFFFFF", "CBD5E1")
                badge = rect(slide, x + Inches(0.16), y - Inches(0.2),
                             Inches(0.46), Inches(0.46), ACCENT,
                             shape=MSO_SHAPE.OVAL)
                label(badge, num, size=13, color="FFFFFF")
                textbox(slide, x + Inches(0.16), y + Inches(0.42),
                        cw - Inches(0.32), Inches(0.36), head, size=14.5,
                        bold=True, color=ACCENT_DK)
                textbox(slide, x + Inches(0.16), y + Inches(0.84),
                        cw - Inches(0.32), ch - Inches(1.0), body, size=11.5,
                        color=MUTED, spacing=1.1)

        # ── scorecard ──
        elif kind == "scorecard":
            rows_ = spec["rows"]
            # fit however many rows there are into the space below the title bar
            avail = SH - top - Inches(0.55)
            rowh = min(Inches(0.60), int(avail / max(len(rows_), 1)) - Inches(0.04))
            dot = int(min(Inches(0.28), rowh - Inches(0.22)))
            pad = int((rowh - dot) / 2)
            y = top + Inches(0.05)
            for n, (dim, score, note) in enumerate(rows_):
                bg = "F8FAFC" if n % 2 == 0 else "FFFFFF"
                rect(slide, Inches(0.55), y, SW - Inches(1.1), rowh, bg,
                     shape=MSO_SHAPE.RECTANGLE)
                textbox(slide, Inches(0.78), y + pad - Inches(0.06), Inches(3.1),
                        Inches(0.4), dim, size=14, bold=True, color=INK)
                col = GOOD if score >= 4 else (WARN if score == 3 else BAD)
                for s in range(5):
                    filled = s < score
                    rect(slide, Inches(4.0) + s * Inches(0.36),
                         y + pad, dot, dot,
                         col if filled else "E2E8F0", shape=MSO_SHAPE.OVAL)
                textbox(slide, Inches(6.15), y + pad - Inches(0.05),
                        SW - Inches(6.8), Inches(0.4), note, size=12.5,
                        color=MUTED)
                y += rowh + Inches(0.04)

        # ── roadmap ──
        elif kind == "roadmap":
            phases = spec["phases"]
            cw = (SW - Inches(1.5) - Inches(0.26) * 3) / 4
            colors_ = [(BAD, "FEF2F2"), (WARN, "FFFBEB"), (ACCENT, "EFF6FF"),
                       (GOOD, "F0FDF4")]
            for n, (name, dur, head, body) in enumerate(phases):
                x = Inches(0.75) + n * (cw + Inches(0.26))
                stroke, fill = colors_[n]
                card = rect(slide, x, top + Inches(0.35), cw, Inches(3.9),
                            fill, stroke)
                cap = rect(slide, x, top + Inches(0.35), cw, Inches(0.62),
                           stroke, shape=MSO_SHAPE.RECTANGLE)
                label(cap, f"{name}   {dur}", size=13.5, color="FFFFFF")
                textbox(slide, x + Inches(0.2), top + Inches(1.14),
                        cw - Inches(0.4), Inches(0.62), head, size=16,
                        bold=True, color=stroke)
                textbox(slide, x + Inches(0.2), top + Inches(1.86),
                        cw - Inches(0.4), Inches(2.2),
                        body.replace(" · ", "\n"), size=12.5, color=INK,
                        spacing=1.35)
            textbox(slide, Inches(0.75), SH - Inches(1.05), SW - Inches(1.5),
                    Inches(0.5),
                    "Roughly ten days of work moves operational maturity from "
                    "two stars to four.",
                    size=13, bold=True, color=ACCENT_DK, align=PP_ALIGN.CENTER)

        # ── closing ──
        elif kind == "closing":
            bar = rect(slide, Inches(0.75), top + Inches(0.2),
                       SW - Inches(1.5), Inches(0.95), "EEF3FB", ACCENT)
            label(bar, spec["lead"], size=21, color=ACCENT_DK)
            textbox(slide, Inches(1.1), top + Inches(1.55), SW - Inches(2.2),
                    Inches(3.4), spec["body"], size=15.5, color=INK,
                    spacing=1.4)

    prs.save(path)


# ─────────────────────────────────────────────────────────────── main ──

def main() -> int:
    print("=" * 66)
    print(f"  {TITLE} — documentation build")
    print("=" * 66)

    missing = [f for f, _ in SOURCES if not os.path.exists(os.path.join(HERE, f))]
    if missing:
        print("  ERROR: missing source documents: " + ", ".join(missing))
        return 1

    total = 0
    for fname, _ in SOURCES:
        size = os.path.getsize(os.path.join(HERE, fname))
        total += size
        print(f"  source   {fname:<38} {size / 1024:8.1f} KB")
    print(f"  {'':<49}{'-' * 11}")
    print(f"  {'total source':<49}{total / 1024:8.1f} KB\n")

    outputs = []

    md_path = os.path.join(HERE, "Architecture_Guide.md")
    with open(md_path, "w", encoding="utf-8") as fh:
        fh.write(build_master_md())
    outputs.append(md_path)
    print(f"  [1/4] Markdown  -> Architecture_Guide.md")

    docx_path = os.path.join(HERE, "Architecture_Guide.docx")
    build_docx(docx_path)
    outputs.append(docx_path)
    print(f"  [2/4] Word      -> Architecture_Guide.docx")

    pdf_path = os.path.join(HERE, "Architecture_Guide.pdf")
    build_pdf(pdf_path)
    outputs.append(pdf_path)
    print(f"  [3/4] PDF       -> Architecture_Guide.pdf")

    pptx_path = os.path.join(HERE, "Architecture_Overview.pptx")
    build_pptx(pptx_path)
    outputs.append(pptx_path)
    print(f"  [4/4] PowerPoint-> Architecture_Overview.pptx")

    print("\n" + "-" * 66)
    ok = True
    for p in outputs:
        if os.path.exists(p) and os.path.getsize(p) > 2048:
            print(f"  OK   {os.path.basename(p):<32} {os.path.getsize(p) / 1024:8.1f} KB")
        else:
            print(f"  FAIL {os.path.basename(p)}")
            ok = False
    print("-" * 66)
    print("  Build complete." if ok else "  Build FAILED.")
    return 0 if ok else 1


if __name__ == "__main__":
    sys.exit(main())
